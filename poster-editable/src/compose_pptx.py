# -*- coding: utf-8 -*-
"""layout.py 를 미리캔버스용 PPTX 로도 굽는다.

PDF 와 같은 원칙이다 — 도형/아이콘/배경은 전부 개별 '그림' 개체,
글자는 편집 가능한 텍스트 상자. PPTX 는 개체 모델 자체가 개별 도형 목록이라
임포터가 한 덩어리로 묶을 여지가 없다.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Mm, Pt
from PIL import Image

import layout as L

ASSETS = r"C:\Claude\poster-editable\assets"
OUT = r"C:\Claude\poster-editable\out"

# Noto Sans KR 세로 메트릭 (em 기준) — 텍스트 상자 위쪽에서 베이스라인까지의 거리
ASCENT = 1.160
DESCENT = 0.288
PT2MM = 25.4 / 72.0

FAMILY = {"NotoR": ("Noto Sans KR", False),
          "NotoM": ("Noto Sans KR Medium", False),
          "NotoB": ("Noto Sans KR", True)}


def set_spacing(run, track_em, size_pt):
    """자간(1/100 pt 단위)을 rPr 에 직접 설정한다 — python-pptx 에 API 가 없다."""
    spc = int(round(track_em * size_pt * 100))
    if spc:
        run.font._rPr.set("spc", str(spc))


def add_runs(slide, op, rs):
    size = op["size"]
    line_h = (ASCENT + DESCENT) * size * PT2MM          # mm
    top = op["y"] - ASCENT * size * PT2MM               # 베이스라인 -> 상자 위쪽
    align = op["align"]

    box_w = 170.0
    if align == "c":
        left, pa = op["x"] - box_w / 2.0, PP_ALIGN.CENTER
    elif align == "r":
        left, pa = op["x"] - box_w, PP_ALIGN.RIGHT
    else:
        left, pa = op["x"], PP_ALIGN.LEFT

    tb = slide.shapes.add_textbox(Mm(left), Mm(top), Mm(box_w), Mm(line_h * 1.25))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = pa
    name, bold = FAMILY[op["font"]]
    for s, color in rs:
        r = p.add_run()
        r.text = s
        r.font.name = name
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(*color)
        set_spacing(r, op["track"], size)
    return tb


def add_img(slide, op):
    p = os.path.join(ASSETS, op["f"])
    if not os.path.exists(p):
        print("  !! MISSING ASSET:", op["f"])
        return
    nw, nh = Image.open(p).size
    w, h = op["w"], op["h"]
    if w is None and h is None:
        w = nw / 8.0
    if w is None:
        w = h * nw / nh
    if h is None:
        h = w * nh / nw
    x, y = op["x"], op["y"]
    if op["anchor"] == "c":
        x -= w / 2.0; y -= h / 2.0
    elif op["anchor"] == "tc":
        x -= w / 2.0
    return slide.shapes.add_picture(p, Mm(x), Mm(y), Mm(w), Mm(h))


def build(path=None):
    prs = Presentation()
    prs.slide_width = Mm(L.PAGE_W)
    prs.slide_height = Mm(L.PAGE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])      # 빈 레이아웃

    n_img = n_txt = 0
    for op in L.build():
        if op["kind"] == "img":
            add_img(slide, op); n_img += 1
        elif op["kind"] == "runs":
            add_runs(slide, op, op["rs"]); n_txt += 1
        else:
            add_runs(slide, op, [(op["s"], op["color"])]); n_txt += 1

    out = path or os.path.join(OUT, "poster-editable.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print("pptx -> %s   그림=%d  텍스트=%d  %.0f KB"
          % (out, n_img, n_txt, os.path.getsize(out) / 1024))
    return out


if __name__ == "__main__":
    build()
